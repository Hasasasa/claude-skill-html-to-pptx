"""assemble.py — 把 measurement JSON 装配成 pptx。

Usage:
    python assemble.py <measurement.json> <out.pptx>

设计：
- 标准 16:9 幻灯片 = 13.333" × 7.5" = 12192000 × 6858000 EMU
- 测量视口 1920×1080 → 1 CSS px = 6350 EMU = 0.5 pt
- pptx 内部直接走低层 lxml 操作 spPr / txBody，避开 python-pptx 高层 API 的限制
"""
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

SLIDE_W_PX = 1920
SLIDE_H_PX = 1080
SLIDE_W_EMU = 12192000      # 13.333"
SLIDE_H_EMU = 6858000       # 7.5"
PX_TO_EMU = SLIDE_W_EMU / SLIDE_W_PX  # 6350
PX_TO_PT = 0.5

from embed_fonts import (
    family_alias_map, cjk_typefaces, cjk_for_style, style_of_typeface,
)
from text_utils import is_cjk_text

# 字体映射全部从 embed_fonts.FONT_PLAN 派生；不再在这里维护单独的表
# - FONT_FALLBACKS: CSS 名 → OOXML typeface
# - CJK_FONTS:      标记为 CJK 的 typeface 集合
# FONT_PLAN 是运行时填充的（font_resolver 按需解析），所以 convert.py 在 resolve
# 之后必须调 refresh_font_plan_caches() 把这三个 module-level 缓存重新派生一次。
FONT_FALLBACKS: dict[str, str] = {}
CJK_FONTS: set[str] = set()
_CJK_ALIAS_SET: set[str] = set()


def refresh_font_plan_caches():
    """font_resolver 改完 FONT_PLAN 之后调一次，让 first_font / cjk_font 看到新条目。"""
    global FONT_FALLBACKS, CJK_FONTS, _CJK_ALIAS_SET
    FONT_FALLBACKS = family_alias_map()
    CJK_FONTS = cjk_typefaces()
    _CJK_ALIAS_SET = {name.lower() for name, tf in FONT_FALLBACKS.items() if tf in CJK_FONTS}


# import 时跑一次（FONT_PLAN 可能此时已被预填，例如 embed CLI 单独跑）
refresh_font_plan_caches()

def parse_text_shadow(value: str):
    """解析 CSS text-shadow，返回 (dx_px, dy_px, blur_px, (r,g,b,a)) 或 None。
    多层 shadow 取第一层。值如：
        "rgba(229, 57, 42, 1) 5px 5px 0px"
        "5px 5px rgb(0, 0, 0)"
    """
    if not value or value == "none":
        return None
    first = value.split(",", 1)[0]
    # 把可能在前后的 rgb()/rgba() 抠出来
    rgba_m = re.search(r"rgba?\(([^)]+)\)", first)
    color_rgba = (0, 0, 0, 1.0)
    if rgba_m:
        parts = [p.strip() for p in rgba_m.group(1).split(",")]
        if len(parts) >= 3:
            color_rgba = (int(float(parts[0])), int(float(parts[1])), int(float(parts[2])),
                          float(parts[3]) if len(parts) >= 4 else 1.0)
        first = re.sub(r"rgba?\([^)]+\)", "", first)
    nums = [float(m.group(1)) for m in re.finditer(r"(-?\d+\.?\d*)px", first)]
    if len(nums) < 2:
        return None
    dx, dy = nums[0], nums[1]
    blur = nums[2] if len(nums) >= 3 else 0.0
    return (dx, dy, blur, color_rgba)


def parse_rgb(s: str):
    """返回 (r,g,b) 元组；丢弃 alpha。要 alpha 请用 parse_rgba。"""
    return parse_rgba(s)[:3]


def parse_rgba(s: str):
    """返回 (r,g,b,a) 其中 a 是 0.0–1.0 浮点；缺省 1.0。"""
    m = re.match(r"rgba?\(([^)]+)\)", s)
    if not m:
        return (0, 0, 0, 1.0)
    parts = [p.strip() for p in m.group(1).split(",")]
    r = int(float(parts[0])); g = int(float(parts[1])); b = int(float(parts[2]))
    a = float(parts[3]) if len(parts) >= 4 else 1.0
    return (r, g, b, a)


GENERIC_FONT_KEYWORDS = {
    # CSS 通用字体族关键字，不算"具体字体"，遇到要跳过
    "serif", "sans-serif", "monospace", "cursive", "fantasy",
    "system-ui", "ui-serif", "ui-sans-serif", "ui-monospace",
    "math", "emoji", "fangsong",
}


def first_font(font_family: str) -> str:
    """从 CSS font-family 字符串里挑第一项（latin 用），去引号。
    优先返回 FONT_FALLBACKS 里映射到的 OOXML typeface；跳过 generic 关键字。"""
    items = [x.strip().strip('"').strip("'") for x in font_family.split(",")]
    for it in items:
        if not it or it.lower() in GENERIC_FONT_KEYWORDS:
            continue
        # 大小写不敏感查 FONT_FALLBACKS（处理 CSS 大小写差异）
        if it in FONT_FALLBACKS:
            return FONT_FALLBACKS[it]
        if it.lower() in FONT_FALLBACKS:
            return FONT_FALLBACKS[it.lower()]
        return it  # 用户用了我们没装的字体，原名透传（运行时回退到系统）
    return items[0] if items else "Calibri"


def cjk_font(font_family: str, latin_name: str) -> str:
    """返回该 run 应使用的 East Asian 字体。

    决策顺序（全部从 FONT_PLAN 派生，无硬编码）：
    1. CSS family 列表里显式列了 CJK 字体（或其 alias） → 用映射后的 typeface
    2. 否则按 latin 字体的 style 配对：latin serif → CJK serif，latin sans/mono → CJK sans
    3. style 未知时默认 sans CJK
    """
    items = [x.strip().strip('"').strip("'") for x in font_family.split(",")]
    for it in items:
        if it.lower() in _CJK_ALIAS_SET:
            return FONT_FALLBACKS.get(it, FONT_FALLBACKS.get(it.lower(), it))
    # 没有显式 CJK：按 latin 字体的 style 选配对的 CJK
    latin_style = style_of_typeface(latin_name)
    return cjk_for_style(latin_style)


def px_to_emu(px: float) -> int:
    return int(round(px * PX_TO_EMU))


def _find_blank_layout(prs):
    """从 prs.slide_layouts 找空白布局，避免硬编码下标 6。

    顺序：
    1. layout.name == 'Blank'（python-pptx 默认模板里就叫这个名字）
    2. 没有 placeholder 的 layout（结构上等同空白）
    3. 兜底用最后一个 layout
    """
    layouts = list(prs.slide_layouts)
    for layout in layouts:
        if (layout.name or "").strip().lower() == "blank":
            return layout
    for layout in layouts:
        if len(list(layout.placeholders)) == 0:
            return layout
    return layouts[-1]


def _text_max_font_size(rec) -> float:
    runs = rec.get("runs", []) or []
    return max((float(run.get("fontSize", 16) or 16) for run in runs), default=16)


def _has_explicit_break(rec) -> bool:
    """rec 的 runs 里是否有显式 `<br>` / `\n` 分行。"""
    for run in (rec.get("runs", []) or []):
        text = run.get("text", "") or ""
        if run.get("linebreak") or (text.strip() and "\n" in text):
            return True
    return False


def _text_is_single_line(rec, max_fs: float) -> bool:
    """判断 record 是不是单行文本。
    用 BCR.h 减掉 CSS padding（垂直方向）得到内容高度，再跟 max_fs*1.8 比。
    不扣 padding 的话，带 padding 的短标签（button / badge / pill）会被误判成多行，
    走 wrap=square 让 PPT 度量稍宽时把短词切开（如 "SELECT" → "SELEC"/"T"）。
    """
    r = rec["rect"]
    style = rec.get("style", {}) or {}
    pad_v = (style.get("paddingTop", 0) or 0) + (style.get("paddingBottom", 0) or 0)
    content_h = max(0, r["h"] - pad_v)
    return content_h < max_fs * 1.8 and not _has_explicit_break(rec)


def _text_box_size_px(rec, max_fs: float, is_single_line: bool) -> tuple[float, float]:
    """textbox 几何严格 = HTML BCR width，避免几何外扩破坏 algn=ctr 视觉居中。
    高度给一点宽裕，避免 PPT 行高比浏览器略高时底部被裁。
    PPT 文字度量比浏览器稍宽的"防换行"靠 wrap='none' (单行) + tf.margin (内 padding 吸收)。
    """
    r = rec["rect"]
    style = rec.get("style") or {}
    display = (style.get("display") or "").lower()
    align_items = (style.get("alignItems") or "").lower()
    # flex/grid 容器有 align-items:center/end 时，BCR 实际是容器 BCR（不是文字 BCR）。
    # assemble 这边按 anchor=ctr/b 在 PPT 里做垂直定位 —— cy 必须严格 = r.h，
    # 不能再 h*1.4 撑高，否则文字会被推到容器底部以下（用户看着就是"居中字怎么靠下了"）。
    if ("flex" in display or "grid" in display) and align_items in ("center", "flex-end", "end"):
        return r["w"], r["h"]
    # 显式 <br> / \n 分行：浏览器已经精确测出多行 BCR，cy 严格 = r.h。
    # 再 1.3x 撑高会把下方相邻段盖住（slide 17 链接段贴到上段"关键是..."同位置就是这个 bug）。
    if _has_explicit_break(rec):
        return r["w"], r["h"]
    if is_single_line:
        return r["w"], max(r["h"] * 1.4, max_fs * 1.6)
    return r["w"], r["h"] * 1.3


def _prepare_text_layouts(records):
    """缓存每个 text record 的 layout 计算（max_fs / is_single_line / textbox 尺寸）。
    一次计算两次复用：text_box_overlap_warnings 在 self_check 也会读 _pptx_text_layout。"""
    text_records = [rec for rec in records if rec.get("kind") == "text" and rec.get("rect")]
    for rec in text_records:
        max_fs = _text_max_font_size(rec)
        is_single_line = _text_is_single_line(rec, max_fs)
        w_px, h_px = _text_box_size_px(rec, max_fs, is_single_line)
        rec["_pptx_text_layout"] = {
            "max_fs": max_fs,
            "is_single_line": is_single_line,
            "w_px": w_px,
            "h_px": h_px,
        }


def make_rgb(rgb):
    r, g, b = rgb
    return RGBColor(r, g, b)


def _apply_rotation(shape, rec):
    """如果 record 上有非零 rotation，重写 shape 的 xfrm：
    - off/ext 用 naturalSize（未旋转的元素本尺寸）+ AABB 中心点（rect 是浏览器返回的旋转后 AABB）
    - 加 rot 属性（1/60000 度）
    会让 PPT 按"先放在 AABB 中心、再旋转"的方式还原 HTML 的视觉效果。
    """
    rot_deg = rec.get("rotation") or 0.0
    if abs(rot_deg) < 0.5:
        return
    nat = rec.get("naturalSize") or {}
    nat_w = float(nat.get("w") or 0)
    nat_h = float(nat.get("h") or 0)
    if nat_w <= 0 or nat_h <= 0:
        return
    rect = rec["rect"]
    cx_px = rect["x"] + rect["w"] / 2.0
    cy_px = rect["y"] + rect["h"] / 2.0
    new_x = px_to_emu(cx_px - nat_w / 2.0)
    new_y = px_to_emu(cy_px - nat_h / 2.0)
    new_w = px_to_emu(nat_w)
    new_h = px_to_emu(nat_h)
    spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        return
    xfrm = spPr.find(qn("a:xfrm"))
    if xfrm is None:
        return
    # OOXML rot 单位 = 1/60000 度，正数 = 顺时针
    rot_units = int(round(rot_deg * 60000)) % (360 * 60000)
    xfrm.set("rot", str(rot_units))
    off = xfrm.find(qn("a:off"))
    ext = xfrm.find(qn("a:ext"))
    if off is not None:
        off.set("x", str(new_x))
        off.set("y", str(new_y))
    if ext is not None:
        ext.set("cx", str(new_w))
        ext.set("cy", str(new_h))


def add_background(slide, rgb):
    """整页底色：插入一个全屏 rectangle。"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W_EMU, SLIDE_H_EMU)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = make_rgb(rgb)
    # 把它放到最底层
    sp = shape._element
    spTree = sp.getparent()
    spTree.remove(sp)
    spTree.insert(2, sp)  # 0,1 是 nvGrpSpPr / grpSpPr
    return shape


def add_shape_box(slide, rec):
    """带 background / border 的非文本装饰节点。

    border 处理：每条边单独判断。
    - 全 4 边都有：画带边的矩形
    - 单侧或不对称：用 connectorStraightLine 单独画每条出现的边线
    """
    r = rec["rect"]
    x, y, w, h = px_to_emu(r["x"]), px_to_emu(r["y"]), px_to_emu(r["w"]), px_to_emu(r["h"])
    if w <= 0 or h <= 0:
        return
    deco = rec.get("deco", {})

    sides = {
        "top":    (deco.get("borderTop"),    deco.get("borderTopWidth", 0)),
        "bottom": (deco.get("borderBottom"), deco.get("borderBottomWidth", 0)),
        "left":   (deco.get("borderLeft"),   deco.get("borderLeftWidth", 0)),
        "right":  (deco.get("borderRight"),  deco.get("borderRightWidth", 0)),
    }
    active_sides = [k for k, (present, _) in sides.items() if present]

    # 圆形检测：CSS border-radius >= 50% 或 >= min(w,h)/2 → 当作椭圆/圆
    is_oval = _is_oval(deco.get("borderRadius", ""), r["w"], r["h"])
    prst = MSO_SHAPE.OVAL if is_oval else MSO_SHAPE.RECTANGLE

    # 如果有填充色：画形状（不带 border，border 单独画线）
    fill_shape = None
    if deco.get("hasBg"):
        r_, g_, b_, a_ = parse_rgba(deco["bg"])
        shape = slide.shapes.add_shape(prst, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(r_, g_, b_)
        # 把 CSS alpha 转成 OOXML 的 a:alpha（单位 1/1000 百分比，100000 = 100%）
        if a_ < 1.0:
            _set_fill_alpha(shape, a_)
        shape.line.fill.background()
        _apply_rotation(shape, rec)
        fill_shape = shape

    # 4 边都有边框 → 用形状带 border 一次画完
    if len(active_sides) == 4 and not deco.get("hasBg"):
        shape = slide.shapes.add_shape(prst, x, y, w, h)
        shape.fill.background()
        b_rgba = parse_rgba(deco.get("borderColor", "rgb(127,127,127)"))
        shape.line.color.rgb = make_rgb(b_rgba[:3])
        widest = max(deco.get("borderTopWidth", 0), deco.get("borderBottomWidth", 0),
                     deco.get("borderLeftWidth", 0), deco.get("borderRightWidth", 0))
        shape.line.width = Emu(px_to_emu(widest))
        _set_line_alpha(shape, b_rgba[3])
        _apply_rotation(shape, rec)
        return shape

    # 椭圆形（仅有填充无 4 边边框）：完了直接 return，跳过下面"按需画线"
    if is_oval and deco.get("hasBg"):
        return fill_shape

    # 否则按需画线（border-top / border-bottom 等单侧情形）
    border_rgba = parse_rgba(deco.get("borderColor", "rgb(127,127,127)"))
    rgb = border_rgba[:3]
    alpha = border_rgba[3]
    for side in active_sides:
        bw = sides[side][1] or 1
        if side == "top":
            _add_line(slide, x, y, x + w, y, rgb, bw, alpha)
        elif side == "bottom":
            _add_line(slide, x, y + h, x + w, y + h, rgb, bw, alpha)
        elif side == "left":
            _add_line(slide, x, y, x, y + h, rgb, bw, alpha)
        elif side == "right":
            _add_line(slide, x + w, y, x + w, y + h, rgb, bw, alpha)


def _is_oval(border_radius: str, w_px: float, h_px: float) -> bool:
    """判断元素是否应渲染为椭圆/圆。
    标准：border-radius 是 "50%" 或更大百分比，或者 px 值 >= min(w,h)/2 * 0.9。
    """
    if not border_radius or border_radius == "0px":
        return False
    s = str(border_radius).strip()
    if s.endswith("%"):
        try:
            return float(s[:-1]) >= 50.0   # 去掉最后 1 个 "%" 字符，不是 2
        except ValueError:
            return False
    if s.endswith("px"):
        try:
            px = float(s[:-2])
            # 至少要达到 短边/2 × 0.9 才算椭圆 / pill 形（防止小圆角误判）
            return px >= min(w_px, h_px) / 2 * 0.9
        except ValueError:
            return False
    return False


def _set_fill_alpha(shape, alpha: float):
    """在 shape 的 solidFill 上加 a:alpha 子元素。alpha ∈ [0,1]。"""
    spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        return
    solidFill = spPr.find(qn("a:solidFill"))
    if solidFill is None:
        return
    srgb = solidFill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    # OOXML a:alpha val 单位为 千分之一（100000 = 100%）
    alpha_el = etree.SubElement(srgb, qn("a:alpha"))
    alpha_el.set("val", str(int(round(alpha * 100000))))


def _set_line_alpha(shape_or_line, alpha: float):
    """在 spPr/a:ln/a:solidFill/a:srgbClr 上加 a:alpha 子元素。alpha ∈ [0,1]。
    alpha >= 1.0 时直接返回（OOXML 默认就是不透明）。"""
    if alpha >= 1.0:
        return
    spPr = shape_or_line._element.find(qn("p:spPr"))
    if spPr is None:
        return
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        return
    solid = ln.find(qn("a:solidFill"))
    if solid is None:
        return
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return
    alpha_el = etree.SubElement(srgb, qn("a:alpha"))
    alpha_el.set("val", str(int(round(alpha * 100000))))


def _add_line(slide, x1, y1, x2, y2, color_rgb, width_px, alpha: float = 1.0):
    """画一条直线。x1/y1/x2/y2 已是 EMU。可选 alpha 支持半透明。"""
    from pptx.enum.shapes import MSO_CONNECTOR
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = make_rgb(color_rgb)
    line.line.width = Emu(px_to_emu(width_px))
    _set_line_alpha(line, alpha)
    return line


def _normalize_runs(runs):
    """把测到的 runs 折叠成 OOXML 段落生成器要消费的 (kind, ...) 项序列。

    复刻 CSS 空白折叠：每个 run 内 \\s+ → 单空格，去掉首尾空白，过滤纯空白 run；
    跨 run 的"边界空格"用单独的 SPACE 项保留，避免 run 拼接后丢空格。

    返回 list[(kind, ...)], 其中 kind ∈ {"BREAK", "SPACE", "RUN"}：
    - ("BREAK", run)        显式 <br> / \\n
    - ("SPACE", None)       run 之间的边界空格
    - ("RUN", run, body)    实际文本 run
    """
    cleaned = []
    for i, run in enumerate(runs):
        text = run.get("text", "")
        if not text:
            continue
        if run.get("linebreak"):
            cleaned.append(("BREAK", run))
            continue
        collapsed = re.sub(r"\s+", " ", text)
        leading = collapsed.startswith(" ")
        trailing = collapsed.endswith(" ")
        body = collapsed.strip()
        if not body:
            # 纯空白 run：若非首末位，当作 run 间空格信号
            if cleaned and i != len(runs) - 1:
                cleaned.append(("SPACE", None))
            continue
        if leading and cleaned and cleaned[-1][0] != "BREAK":
            cleaned.append(("SPACE", None))
        cleaned.append(("RUN", run, body))
        if trailing and i != len(runs) - 1:
            cleaned.append(("SPACE", None))

    # 去掉相邻的重复 SPACE
    deduped = []
    for item in cleaned:
        if item[0] == "SPACE" and deduped and deduped[-1][0] == "SPACE":
            continue
        deduped.append(item)
    return deduped


def _fix_empty_paragraph_sizes(tf, style_font_size_px):
    """把空段（<a:p> 无 <a:r>）的 endParaRPr.sz 设成正文字号。

    不显式设的话，OOXML 空段会用 PPT 默认 18pt × 行距撑出 ≈43px 空行高，
    把后续内容推出 textbox 砸到下方相邻段。
    """
    fs_px = float(style_font_size_px or 16) or 16
    end_sz = max(100, int(round(fs_px * 0.5 * 100)))  # px → pt = px*0.5; OOXML sz=pt*100
    for para in tf.paragraphs:
        if para._p.find(qn("a:r")) is None:
            endR = para._p.find(qn("a:endParaRPr"))
            if endR is None:
                endR = etree.SubElement(para._p, qn("a:endParaRPr"))
            endR.set("sz", str(end_sz))


def add_text_box(slide, rec):
    """文本节点 → pptx textbox（多 run 富文本）。"""
    r = rec["rect"]
    # 若 text leaf 同时带背景色 / 边框（如 .bar.a 既是柱子又装着文字），
    # 先按 rec 的 deco 画一个 shape 垫底，再画文字框 —— 否则背景丢失。
    deco = rec.get("deco", {})
    has_decoration = deco.get("hasBg") or deco.get("borderTop") or deco.get("borderBottom") \
                     or deco.get("borderLeft") or deco.get("borderRight")
    if has_decoration:
        # 用一个仅含 deco 的合成 record 调 add_shape_box
        synth = {"rect": r, "deco": deco, "kind": "shape", "tag": rec.get("tag", "div")}
        add_shape_box(slide, synth)

    # 判断单行 / 多行：浏览器测得高度 < 1.8 × 主字号 → 单行
    layout = rec.get("_pptx_text_layout")
    if layout is None:
        max_fs = _text_max_font_size(rec)
        is_single_line = _text_is_single_line(rec, max_fs)
        w_px, h_px = _text_box_size_px(rec, max_fs, is_single_line)
    else:
        max_fs = layout["max_fs"]
        is_single_line = layout["is_single_line"]
        w_px = layout["w_px"]
        h_px = layout["h_px"]

    x = px_to_emu(r["x"])
    y = px_to_emu(r["y"])
    w = px_to_emu(w_px)
    h = px_to_emu(h_px)
    if w <= 0:
        w = px_to_emu(50)
    if h <= 0:
        h = px_to_emu(20)

    tb = slide.shapes.add_textbox(x, y, w, h)
    _apply_rotation(tb, rec)
    tf = tb.text_frame
    # 何时禁用自动 wrap：
    # (1) 真单行（h < max_fs*1.8）— 短标签 PPT 度量稍宽时不应换行
    # (2) 有显式 <br> — 作者已用 <br> 决定分行，PPT 不该在每段里再自动切
    # 自由流动的多行段落（无 <br>，纯靠 word-wrap）保持 wrap=square 让 PPT 按宽度切。
    runs_raw = rec.get("runs") or []
    has_explicit_break = any(run.get("linebreak") for run in runs_raw)
    no_auto_wrap = is_single_line or has_explicit_break
    tf.word_wrap = not no_auto_wrap
    # textbox 几何 = HTML BCR；CSS padding 进 OOXML 内 margin。
    # 这样：text-align:center 时居中点真的对齐 HTML 元素中心；
    #      padding-left 给 ::before marker 留位的设计在 PPT 也保持文字位置。
    style_padding = rec.get("style", {})
    tf.margin_left = px_to_emu(style_padding.get("paddingLeft", 0) or 0)
    tf.margin_right = px_to_emu(style_padding.get("paddingRight", 0) or 0)
    tf.margin_top = px_to_emu(style_padding.get("paddingTop", 0) or 0)
    tf.margin_bottom = px_to_emu(style_padding.get("paddingBottom", 0) or 0)
    # OOXML wrap 属性：square = 框内 wrap，none = 不 wrap 允许溢出
    # 必须与 tf.word_wrap 保持一致，否则 wrap 属性会覆盖 word_wrap 设置
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    bodyPr.set("wrap", "none" if no_auto_wrap else "square")
    # 垂直 anchor：HTML 用 flex/grid + align-items 居中时，OOXML 用 anchor 翻译
    style_for_anchor = rec.get("style", {})
    display = (style_for_anchor.get("display") or "").lower()
    align_items = (style_for_anchor.get("alignItems") or "").lower()
    anchor = "t"
    if "flex" in display or "grid" in display:
        if align_items == "center":
            anchor = "ctr"
        elif align_items in ("flex-end", "end"):
            anchor = "b"
    bodyPr.set("anchor", anchor)
    # 不要 autofit（防止 PPT 自己缩字）
    for child in list(bodyPr):
        if child.tag.endswith("normAutofit") or child.tag.endswith("spAutoFit"):
            bodyPr.remove(child)

    # 处理 runs
    runs = rec.get("runs", [])
    text_transform = rec.get("style", {}).get("textTransform", "none")

    # 第一段
    p = tf.paragraphs[0]
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    # 水平对齐：HTML flex/grid 容器靠 justify-content 居中，普通块靠 text-align。
    # OOXML 段落 algn 二选一：flex/grid 时 justify-content 优先（覆盖 text-align）。
    style_for_algn = rec.get("style", {})
    align = style_for_algn.get("textAlign", "start")
    align_map = {"start": "l", "left": "l", "center": "ctr", "right": "r", "end": "r"}
    justify_map = {"center": "ctr", "flex-end": "r", "end": "r", "right": "r",
                   "flex-start": "l", "start": "l", "left": "l",
                   "space-between": "just", "space-around": "ctr", "space-evenly": "ctr"}
    h_algn = align_map.get(align, "l")
    if "flex" in display or "grid" in display:
        jc = (style_for_algn.get("justifyContent") or "").lower()
        if jc in justify_map:
            h_algn = justify_map[jc]
    pPr.set("algn", h_algn)
    # 显式写行距：根据 HTML 实测 line-height/font-size 比率
    # 不写的话 PPT 用字体默认（CJK 字体默认行距大幅大于 CSS），导致大标题叠压下面元素
    style_for_lh = rec.get("style", {})
    _apply_line_spacing(pPr, style_for_lh.get("lineHeight"), style_for_lh.get("fontSize", 16))

    # 清掉默认 run
    for r_el in p._p.findall(qn("a:r")):
        p._p.remove(r_el)

    first_para = p
    cur_para = first_para

    deduped = _normalize_runs(runs)

    pending_space = False
    for item in deduped:
        kind = item[0]
        if kind == "BREAK":
            cur_para = tf.add_paragraph()
            new_pPr = cur_para._p.get_or_add_pPr()
            new_pPr.set("algn", h_algn)
            _apply_line_spacing(new_pPr, style_for_lh.get("lineHeight"), style_for_lh.get("fontSize", 16))
            pending_space = False
        elif kind == "SPACE":
            pending_space = True
        else:  # RUN
            _, run, body = item
            if pending_space:
                body = " " + body
                pending_space = False
            _emit_run(cur_para, body, run, text_transform)

    _fix_empty_paragraph_sizes(tf, style_for_lh.get("fontSize", 16))

    return tb


def _parse_line_height_pct(line_height: str | float | None, font_size_px: float) -> int | None:
    """把 CSS line-height + font-size 解析成 OOXML spcPct 千分比（单位是 1/1000 %）。

    OOXML <a:spcPct val="120000"/> = 120% 行距。

    输入示例：
    - "184.32px" + fs=192 → 96% → 返回 96000
    - "1.2"             → 120% → 返回 120000
    - "120%"            → 120% → 返回 120000
    - "normal"          → 返回 None（不显式写，但调用方可降级到 120000 兜底）
    """
    if not line_height or not font_size_px or font_size_px <= 0:
        return None
    s = str(line_height).strip()
    if not s or s.lower() == "normal":
        return None
    # 百分号
    if s.endswith("%"):
        try:
            return int(float(s[:-1]) * 1000)
        except ValueError:
            return None
    # px 显式
    if s.endswith("px"):
        try:
            px = float(s[:-2])
            return int((px / font_size_px) * 100000)
        except ValueError:
            return None
    # 无单位数字（CSS line-height 倍数）
    try:
        return int(float(s) * 100000)
    except ValueError:
        return None


def _apply_line_spacing(pPr, line_height, font_size_px, default_pct=120000):
    """在 pPr 上写 <a:lnSpc><a:spcPct val="..."/></a:lnSpc>。

    PPT 对 CJK 字体的默认行距比 CSS 大很多（典型 1.8-2.0x），不显式写就会出现
    大标题撑下来覆盖下方元素的"叠压"。这里把 CSS 的 line-height 显式传给 OOXML。

    line-height 是 "normal" 时用 default_pct（120%）兜底——比 PPT 默认紧但
    不会产生叠压。
    """
    pct = _parse_line_height_pct(line_height, font_size_px)
    if pct is None:
        pct = default_pct
    # 清掉旧的 lnSpc（防重复运行）
    for old in pPr.findall(qn("a:lnSpc")):
        pPr.remove(old)
    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
    spcPct.set("val", str(pct))
    # OOXML 子元素顺序：lnSpc 必须放在 pPr 的开头（在 buNone / buChar 之前）
    # pPr 子元素顺序: lnSpc, spcBef, spcAft, ..., defRPr
    # 简单做法：把 lnSpc 移到最前
    pPr.remove(lnSpc)
    pPr.insert(0, lnSpc)


def _emit_run(paragraph, text, run, text_transform):
    """向 paragraph 写入一个富文本 run。直接操作 OOXML 以精确控制 letterSpacing。"""
    if not text:
        return
    # text-transform
    if text_transform == "uppercase":
        text = text.upper()
    elif text_transform == "lowercase":
        text = text.lower()

    # 解析参数
    font_name = first_font(run.get("fontFamily", "Calibri"))
    font_size_px = run.get("fontSize", 16)
    font_size_pt = round(font_size_px * PX_TO_PT, 2)

    weight = run.get("fontWeight", "400")
    try:
        bold = int(weight) >= 600
    except ValueError:
        bold = weight in ("bold", "bolder")
    italic = run.get("fontStyle", "normal") == "italic"
    color_rgb = parse_rgb(run.get("color", "rgb(0,0,0)"))

    # letter-spacing: 字符串如 "3.6px" / "normal"
    ls = run.get("letterSpacing", "normal")
    if isinstance(ls, str) and ls.endswith("px"):
        ls_px = float(ls[:-2])
    else:
        ls_px = 0.0
    # OOXML spc 单位 = 1/100 pt
    spc_units = int(round(ls_px * PX_TO_PT * 100))

    # 中文 vs 英文：用 ea / latin 区分（CJK 范围定义见 text_utils.CJK_RE）
    is_chinese = is_cjk_text(text)

    # 构造 <a:r><a:rPr ...><a:rFont/></a:rPr><a:t>...</a:t></a:r>
    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    r_el = etree.SubElement(paragraph._p, qn("a:r"))
    rPr = etree.SubElement(r_el, qn("a:rPr"))
    rPr.set("lang", "zh-CN" if is_chinese else "en-US")
    rPr.set("sz", str(int(round(font_size_pt * 100))))  # OOXML sz = 1/100 pt
    if bold:
        rPr.set("b", "1")
    if italic:
        rPr.set("i", "1")
    if spc_units:
        rPr.set("spc", str(spc_units))
    rPr.set("dirty", "0")

    # fill color
    solidFill = etree.SubElement(rPr, qn("a:solidFill"))
    srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgb.set("val", "{:02X}{:02X}{:02X}".format(*color_rgb))

    # text-shadow → OOXML outerShdw（必须在 solidFill 之后、latin 之前）
    shadow = parse_text_shadow(run.get("textShadow", "none"))
    if shadow:
        import math
        dx, dy, blur, (sr, sg, sb, sa) = shadow
        dist_px = math.sqrt(dx * dx + dy * dy)
        if dist_px > 0.5 or blur > 0.5:
            angle_deg = math.degrees(math.atan2(dy, dx))
            if angle_deg < 0:
                angle_deg += 360
            effect_lst = etree.SubElement(rPr, qn("a:effectLst"))
            outer = etree.SubElement(effect_lst, qn("a:outerShdw"))
            outer.set("blurRad", str(int(blur * PX_TO_EMU)))
            outer.set("dist", str(int(dist_px * PX_TO_EMU)))
            outer.set("dir", str(int(angle_deg * 60000)))
            outer.set("algn", "ctr")
            outer.set("rotWithShape", "0")
            shd_clr = etree.SubElement(outer, qn("a:srgbClr"))
            shd_clr.set("val", f"{sr:02X}{sg:02X}{sb:02X}")
            if sa < 1.0:
                shd_a = etree.SubElement(shd_clr, qn("a:alpha"))
                shd_a.set("val", str(int(sa * 100000)))

    # font 分离：latin 用 css 第一项；ea 走 CJK 字体（Noto Serif/Sans SC）
    # PPT 会按字符自动用 latin 还是 ea，所以 IBM Plex Mono 文本里的 CJK 字符
    # 自动落到 Noto Sans SC 上，避免变 tofu。
    ea_name = cjk_font(run.get("fontFamily", ""), font_name)
    latin_el = etree.SubElement(rPr, qn("a:latin"))
    latin_el.set("typeface", font_name)
    ea_el = etree.SubElement(rPr, qn("a:ea"))
    ea_el.set("typeface", ea_name)
    cs_el = etree.SubElement(rPr, qn("a:cs"))
    cs_el.set("typeface", font_name)

    t_el = etree.SubElement(r_el, qn("a:t"))
    t_el.text = text


def add_svg_picture(slide, rec):
    """直接用 measure 阶段已经截好的 SVG PNG。"""
    r = rec["rect"]
    if r["w"] <= 0 or r["h"] <= 0:
        return
    png_path = rec.get("screenshot")
    if not png_path or not Path(png_path).exists():
        print(f"  [skip svg] no screenshot for {rec.get('marker', '?')}")
        return
    slide.shapes.add_picture(png_path,
                             px_to_emu(r["x"]),
                             px_to_emu(r["y"]),
                             px_to_emu(r["w"]),
                             px_to_emu(r["h"]))


def add_img_picture(slide, rec):
    """measure 阶段截好的 <img> 元素 PNG（src 可能是 PNG/JPG/SVG/远程 URL）。
    走截图通道而不是直接嵌 src，避免 cross-origin / SVG 不能直嵌进 OOXML 等问题。"""
    r = rec["rect"]
    if r["w"] <= 0 or r["h"] <= 0:
        return
    png_path = rec.get("screenshot")
    if not png_path or not Path(png_path).exists():
        print(f"  [skip img] no screenshot for {rec.get('src', '?')}")
        return
    slide.shapes.add_picture(png_path,
                             px_to_emu(r["x"]),
                             px_to_emu(r["y"]),
                             px_to_emu(r["w"]),
                             px_to_emu(r["h"]))


def add_canvas_picture(slide, rec):
    """canvas 元素（Chart.js / WebGL / 自绘图）→ picture 嵌入。
    measure 阶段在切到目标页后等待 canvas 像素稳定再截图。
    """
    r = rec["rect"]
    if r["w"] <= 0 or r["h"] <= 0:
        return
    png_path = rec.get("screenshot")
    if not png_path or not Path(png_path).exists():
        print(f"  [skip canvas] no screenshot for {rec.get('marker', '?')}")
        return
    pic = slide.shapes.add_picture(png_path,
                                    px_to_emu(r["x"]),
                                    px_to_emu(r["y"]),
                                    px_to_emu(r["w"]),
                                    px_to_emu(r["h"]))
    _apply_rotation(pic, rec)
    return pic


def add_deco_snapshot(slide, rec):
    """装饰元素截图（background-image / box-shadow / 伪元素装饰 / 非平移 transform）
    → picture 嵌入。

    重要：Playwright `locator.screenshot()` 对带 transform 的元素截的是 AABB
    （旋转/skew 后的可见矩形），旋转已经"烘焙"到位图里。直接放在 AABB rect
    (rec.rect.x/y/w/h) 即可——**不再调 `_apply_rotation`**，否则会双重旋转
    + naturalSize 压缩，对大尺寸旋转矩形（全宽 ribbon 等）尤其灾难。

    子节点的文字 / 子装饰仍按原流程被绘制在它之上。
    """
    r = rec["rect"]
    if r["w"] <= 0 or r["h"] <= 0:
        return
    png_path = rec.get("screenshot")
    if not png_path or not Path(png_path).exists():
        return
    pic = slide.shapes.add_picture(png_path,
                                    px_to_emu(r["x"]),
                                    px_to_emu(r["y"]),
                                    px_to_emu(r["w"]),
                                    px_to_emu(r["h"]))
    return pic


def assemble_slide(slide, data):
    """装配一张 slide。"""
    bg_rgb = parse_rgb(data["slide"]["background"])
    add_background(slide, bg_rgb)
    _prepare_text_layouts(data["records"])

    text_records = []
    for rec in data["records"]:
        if rec["kind"] == "shape":
            # 整页 section 不再单独 add（背景已铺）
            if rec["rect"]["w"] >= SLIDE_W_PX * 0.99 and rec["rect"]["h"] >= SLIDE_H_PX * 0.99:
                continue
            add_shape_box(slide, rec)
        elif rec["kind"] == "text":
            text_records.append(rec)
        elif rec["kind"] == "svg":
            add_svg_picture(slide, rec)
        elif rec["kind"] == "canvas":
            add_canvas_picture(slide, rec)
        elif rec["kind"] == "deco_snapshot":
            add_deco_snapshot(slide, rec)
        elif rec["kind"] == "img":
            add_img_picture(slide, rec)

    # Text sits above rasterized SVG/canvas/deco snapshots. Otherwise an opaque
    # picture can cover positioned labels that belong visually on top of it.
    for rec in text_records:
        add_text_box(slide, rec)


def assemble(measurement, out_path: Path):
    """measurement 可以是 dict（in-process 调用）或 Path（CLI 调用）。"""
    if isinstance(measurement, (str, Path)):
        data = json.loads(Path(measurement).read_text(encoding="utf-8"))
    else:
        data = measurement

    # 兼容单页和多页
    if "slides" in data:
        slides_data = data["slides"]
    else:
        slides_data = [data]

    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU
    blank_layout = _find_blank_layout(prs)

    for i, sdata in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        assemble_slide(slide, sdata)
        print(f"  page {i+1:02d}: {len(sdata.get('records', []))} records, theme={sdata['slide']['theme']}")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"saved {out_path} ({out_path.stat().st_size:,} B, {len(slides_data)} slides)")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)
    assemble(Path(sys.argv[1]).resolve(), Path(sys.argv[2]).resolve())
